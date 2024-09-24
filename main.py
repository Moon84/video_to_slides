import os
import time
import sys
import cv2
import re
from frame_differencing import capture_slides_frame_diff
from post_process import remove_duplicates
# from utils import resize_image_frame
from datetime import datetime
import argparse

import numpy as np
import json
from video_whisper import video_whisper

# import pytesseract
import easyocr

from pptx import Presentation
from pptx.util import Inches



    
# -------------- Initializations ---------------------

FRAME_BUFFER_HISTORY = 15   # Length of the frame buffer history to model background.   边框？
DEC_THRESH = 0.75           # Threshold value, above which it is marked foreground, else background.
DIST_THRESH = 100           # Threshold on the squared distance between the pixel and the sample to decide whether a pixel is close to that sample.

MIN_PERCENT = 0.15          # %age threshold to check if there is motion across subsequent frames
MAX_PERCENT = 0.01          # %age threshold to determine if the motion across frames has stopped.

# pytesseract.pytesseract

'''
    Opencv support avi,mp4,mpeg ,mkv,flv etc. filetype.
    
'''


# ----------------------------------------------------



def resize_image_frame(frame, resize_width):  ##调整图像的尺寸

    ht, wd, _ = frame.shape
    new_height = resize_width * ht / wd
    frame = cv2.resize(frame, (resize_width, int(new_height)), interpolation=cv2.INTER_AREA)

    return frame


def capture_slides_bg_modeling(video_path, output_dir_path, type_bgsub='GMG', history=15, threshold=0.75, min_percent_thresh=0.15, max_percent_thresh=0.01): 
    """
        使用指定的背景建模方法捕获视频中的幻灯片。
        参数:
            video_path (str): 视频文件的路径。
            output_dir_path (str): 存储捕获幻灯片图像的输出目录。
            type_bgsub (str): 用于背景建模的算法类型，可以是 'GMG' 或 'KNN'。
            history (int): 背景建模算法的历史帧数。
            threshold (float): 用于前景/背景分割的阈值。
            MIN_PERCENT_THRESH (float): 用于判断是否存在运动的百分比阈值。
            MAX_PERCENT_THRESH (float): 用于确定运动是否停止的百分比阈值。
        返回:
            无
        此函数读取视频文件，使用指定的背景建模算法（'GMG' 或 'KNN'）对每一帧进行处理。根据算法输出的前景掩码计算非零像素的百分比，以确定是否存在运动。如果在连续帧中运动停止（低于 MAX_PERCENT_THRESH），则捕获当前帧作为幻灯片并保存为图像文件。函数最后打印统计信息，包括总时间、捕获的幻灯片数量。
    """

    # 打印当前使用的背景建模类型
    print(f"Using {type_bgsub} for Background Modeling...")
    # 打印分隔符
    print('---'*10)
    # 根据选择的背景建模类型初始化对应的背景减除器

    os.makedirs(output_dir_path, exist_ok=True)
    print('Output directory %s created...' % output_dir_path)

    video_name, file_ext = os.path.splitext(os.path.basename(video_path))
    # output_dir_path = os.path.join('project/',video_name,'/out_put' )
    # os.makedirs(output_dir_path, exist_ok=True)
    
    if type_bgsub == 'GMG':
        # 使用 GMG 算法进行背景建模
        bg_sub = cv2.bgsegm.createBackgroundSubtractorGMG(initializationFrames=history, decisionThreshold=threshold)
    elif type_bgsub == 'KNN':
        # 使用 KNN 算法进行背景建模
        bg_sub = cv2.createBackgroundSubtractorKNN(history=history, dist2Threshold=threshold, detectShadows=False) 
        
    # 初始化捕获帧标志为 False，表示没有捕获到幻灯片
    capture_frame = False
    # 初始化捕获的幻灯片数量为 0
    screenshots_count = 0

    # 打开视频文件进行逐帧捕获
    cap = cv2.VideoCapture(video_path)

    # 检查视频文件是否可以成功打开
    if not cap.isOpened():
        # 如果无法打开视频文件，则打印错误信息并退出程序
        print('Unable to open video file: ' % video_path)
        sys.exit()
     
    
    # 记录开始时间
    start = time.time()

    
    # 当视频文件打开时，循环读取后续帧
    while cap.isOpened():
        # 读取一帧图像
        ret, frame = cap.read()
        # 获取视频的帧率和总帧数
        fps = int(cap.get(cv2.CAP_PROP_FPS))
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        current_frame = int(cap.get(cv2.CAP_PROP_POS_FRAMES))

        # 如果读取帧失败或视频结束，则退出循环
        if not ret:
            break

        # 创建原始帧的副本
        orig_frame = frame.copy() 
        # 调整帧大小以保持宽高比
        frame = resize_image_frame(frame, resize_width=640) 

        # 将当前帧应用到背景减除器
        fg_mask = bg_sub.apply(frame) 

        # 计算前景掩码中的非零像素百分比
        p_non_zero = (cv2.countNonZero(fg_mask) / (1.0 * fg_mask.size)) * 100

        # 当非零像素百分比低于 MAX_PERCENT_THRESH 时，表示运动已停止，捕获当前帧
        if p_non_zero < max_percent_thresh and not capture_frame:
            capture_frame = True
            # 将当前帧转换为秒
            time_in_seconds = current_frame / fps

            # 递增截图计数
            screenshots_count += 1
            
            # 将时间转换为格式化的字符串，用于文件名
            png_filename = '{:02d}:{:02d}:{:02d}.jpg'.format(int(time_in_seconds // 3600), int((time_in_seconds % 3600) // 60), int(time_in_seconds % 60))
            # 构建输出文件的完整路径
            out_file_path = os.path.join(output_dir_path, png_filename)
            # 打印保存文件的信息
            print(f"Saving file at: {out_file_path}")
            # 保存当前帧为图像文件
            cv2.imwrite(out_file_path, orig_frame)
            

        # 当捕获帧标志为真且非零像素百分比大于等于 MIN_PERCENT_THRESH 时，等待直到运动在后续帧中停止
        elif capture_frame and p_non_zero >= min_percent_thresh:
            capture_frame = False


    # 记录结束时间
    end_time = time.time()
    # 打印分隔符
    print('***'*10,'\n')
    # 打印统计信息
    print("Statistics:")
    # 打印分隔符
    print('---'*10)
    # 打印总捕获时间，保留3位小数
    print(f'Total Time taken: {round(end_time-start, 3)} secs')
    # 打印总捕获的幻灯片数量
    print(f'Total Screenshots captured: {screenshots_count}')
    # 打印分隔符
    print('---'*10,'\n')
    
    # 释放视频捕获对象，关闭视频文件
    cap.release()



def convert_slides_to_ppt(output_dir_path):   ##把文件夹下的图片按照时间序列排序，转换成PPT
    slide_number = 0  # 初始化 slide_number
    # 创建PPT对象
    ppt = Presentation()
    ppt.slide_width = Inches(16)  # 设置幻灯片宽度为10英寸
    ppt.slide_height = Inches(9)  # 设置幻灯片高度为5.625英寸

    print('Output folder_path is', output_dir_path)

    seq = 0
    pic_list = []
    reader = easyocr.Reader(['en']) # this needs to run only once to load the model into memory        


    # 遍历图片，转化成json格式
    for file_name in sorted(os.listdir(output_dir_path)):    ##按照文件名进行排序后再遍历
        if file_name.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')) :
            image_path = os.path.join(output_dir_path, file_name)
            seq += 1
            timestamp = file_name.split('.')[0]
            ocr_text = reader.readtext(image_path, detail=0) # 识别图片中的文本
            ocr_text = re.split(r'\s+', ' '.join(ocr_text)) # 将所有文本连接后再按空格分割成单词
            # ocr_text = ocr_text.replace('\n','')  # 替换换行符为空格
            # print("ocr_text type is %s" % type(ocr_text))

           # 创建结果字典
            result = {
                "seq": seq,
                "file_path": image_path,
                "timestamp": timestamp,
                "ocr_text": ocr_text,
                "content": '',
            }
            
            pic_list.append(result)
    

    for i in range(len(pic_list) - 1):
        if i==(len(pic_list)-1):
            break

        current_image = pic_list[i]
        next_image = pic_list[i+1]
        
        
        # 检查当前图片的OCR文本是否完全包含在下一张图片中
        if set(current_image['ocr_text']) <= set(next_image['ocr_text']):
            os.remove(current_image['file_path'])
            pic_list.remove(current_image)
            print(f"删除图片: {current_image['file_path']}")
    print('pic_list length is %s' % len(pic_list))


    for i in range(len(pic_list) - 1):
        if i==(len(pic_list)-1):
            print('break at %s'% pic_list[i])
            break
        pic_currnet_json = pic_list[i]
        pic_next_json = pic_list[i+1]

        start_time = datetime.strptime(pic_currnet_json['timestamp'], '%H:%M:%S')
        end_time = datetime.strptime(pic_next_json['timestamp'], '%H:%M:%S')
        relevant_texts = []
        with open('./project/Neural_Networks_Overview/out_put/transcription.json', 'r', encoding='utf-8') as f:
            json_data = json.load(f)
        for trans in json_data:
            trans_start_time = trans['start_time']
            trans_start = datetime.strptime(trans_start_time, '%H:%M:%S.%f')
            trans_start = trans_start.replace(microsecond=0)  # 去掉毫秒部分
            if start_time <= trans_start < end_time:
                relevant_texts.append(trans['text'])
        pic_list[i]['content'] = ' '.join(relevant_texts)
        last_text =  json_data[-1]['text']
        pic_list[i]['content'] = last_text
        print(pic_list[i]['content'])


    for pic in pic_list:    ##按照文件名进行排序后再遍历
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # 添加新的幻灯片
        
        image_path = pic['file_path']
        left, top, width, height = Inches(0.2), Inches(0.2), Inches(9), Inches(9/1280*780)

        slide.shapes.add_picture(image_path, left, top, width=width, height=height)  # 添加图片到幻灯片中

        slide.notes_slide.notes_text_frame.text = pic['content']
        print(image_path,"added to the slides")
        slide_number = slide_number + 1


    ppt_file_name = os.path.basename(video_path).split('.')[0] + '.pptx'
    output_ppt_path = os.path.join(output_dir_path,ppt_file_name)  
    print("output_ppt_path is", output_ppt_path,'\n')
    # 检查文件是否存在
    if not os.path.exists(output_ppt_path):
        # 如果文件不存在，则保存PPT
        ppt.save(output_ppt_path)  # 保存为 output.pptx
        print('PPT Created! under', output_ppt_path)
        print('***'*10,'\n')
    else:
        os.remove(output_ppt_path)
        ppt.save(output_ppt_path)  # 保存为 output.pptx






if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="使用 Whisper 将视频转录为带时间戳的 JSON 文件")
    parser.add_argument("video_path", help="输入视频文件的路径",default='Neural_Networks_Overview.mp4')
    parser.add_argument("-d", "--output_dir_path", default="./project/Neural_Networks_Overview/out_put", help="输出 JSON 文件的路径 (默认: 根文件夹下)")
    args = parser.parse_args()
    video_path = args.video_path
    output_dir_path = args.output_dir_path  
    # video_path = 'Neural_Networks_Overview.mp4'
    # output_dir_path = './project/Neural_Networks_Overview/out_put'
    os.makedirs(output_dir_path, exist_ok=True)
    type_bgsub = 'KNN'
    history = 15
    threshold = 0.75
    min_percent_thresh = 0.15
    max_percent_thresh = 0.01

    capture_slides_bg_modeling(video_path, output_dir_path, type_bgsub, history, threshold, min_percent_thresh, max_percent_thresh)
    remove_duplicates(output_dir_path)
    video_whisper(video_path,output_dir_path,'base')

    convert_slides_to_ppt(output_dir_path)





